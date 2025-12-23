import streamlit as st
import requests
import pandas as pd
import json
import os
import re
import numpy as np
import textwrap
import calendar
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import subprocess, tempfile, zipfile
from datetime import datetime
from io import BytesIO
from pdf2image import convert_from_bytes
from PIL import Image, ImageOps, ImageFilter, ImageChops
from collections import Counter

# ============================================================
#  CONFIGURACIÓN GLOBAL
# ============================================================
MESES_ES = ["enero", "febrero", "marzo", "abril", "mayo", "junio",
            "julio", "agosto", "septiembre", "octubre", "noviembre", "diciembre"]
DIAS_ES = ["lunes", "martes", "miercoles", "jueves", "viernes", "sabado", "domingo"]

# ============================================================
#  FUNCIONES TÉCNICAS
# ============================================================

def recorte_inteligente_bordes(img, umbral_negro=60):
    img_gray = img.convert("L")
    arr = np.array(img_gray)
    h, w = arr.shape
    def fila_es_negra(fila): return (np.sum(fila < 35) / fila.size) * 100 > umbral_negro
    def columna_es_negra(col): return (np.sum(col < 35) / col.size) * 100 > umbral_negro
    top = 0
    while top < h and fila_es_negra(arr[top, :]): top += 1
    bottom = h - 1
    while bottom > top and fila_es_negra(arr[bottom, :]): bottom -= 1
    left = 0
    while left < w and columna_es_negra(arr[:, left]): left += 1
    right = w - 1
    while right > left and columna_es_negra(arr[:, right]): right -= 1
    if right <= left or bottom <= top: return img
    return img.crop((left, top, right + 1, bottom + 1))

def procesar_imagen_inteligente(img_data, target_w_pt, target_h_pt, con_blur=False):
    base_w, base_h = int(target_w_pt / 9525), int(target_h_pt / 9525)
    render_w, render_h = base_w * 2, base_h * 2
    img = Image.open(BytesIO(img_data)).convert("RGB")
    img = recorte_inteligente_bordes(img, umbral_negro=60)
    if con_blur:
        fondo = ImageOps.fit(img, (render_w, render_h), Image.Resampling.LANCZOS)
        fondo = fondo.filter(ImageFilter.GaussianBlur(radius=10))
        img.thumbnail((render_w, render_h), Image.Resampling.LANCZOS)
        offset = ((render_w - img.width) // 2, (render_h - img.height) // 2)
        fondo.paste(img, offset)
        img_final = fondo
    else:
        img_final = img.resize((render_w, render_h), Image.Resampling.LANCZOS)
    output = BytesIO()
    img_final.save(output, format="JPEG", quality=85, subsampling=0, optimize=True)
    output.seek(0)
    return output

def procesar_texto_maestro(texto, campo=""):
    if not texto or str(texto).lower() == "none": return ""
    if isinstance(texto, list): return texto
    if campo == 'Hora': return str(texto).lower().strip()
    t = str(texto).replace('/', ' ').strip().replace('\n', ' ').replace('\r', ' ')
    t = re.sub(r'\s+', ' ', t)
    if campo == 'Seccion': return t.upper()
    palabras = t.lower().split()
    if not palabras: return ""
    prep = ['de', 'la', 'el', 'en', 'y', 'a', 'con', 'las', 'los', 'del', 'al']
    resultado = []
    for i, p in enumerate(palabras):
        es_inicio = (i == 0)
        despues_parentesis = (i > 0 and "(" in palabras[i-1])
        if es_inicio or despues_parentesis or (p not in prep):
            if p.startswith("("): resultado.append("(" + p[1:].capitalize())
            else: resultado.append(p.capitalize())
        else:
            resultado.append(p)
    return " ".join(resultado)

def generar_pdf(pptx_bytes):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        tmp.write(pptx_bytes)
        path = tmp.name
    try:
        subprocess.run(['soffice', '--headless', '--convert-to', 'pdf', '--outdir', os.path.dirname(path), path], check=True)
        pdf_path = path.replace(".pptx", ".pdf")
        with open(pdf_path, "rb") as f: data = f.read()
        os.remove(path); os.remove(pdf_path)
        return data
    except: return None

# --- LÓGICA DE DATOS ---
def obtener_fecha_texto(fecha_dt):
    dia_idx = fecha_dt.weekday()
    return f"{DIAS_ES[dia_idx]} {fecha_dt.day} de {MESES_ES[fecha_dt.month - 1]} de {fecha_dt.year}"

def obtener_hora_texto(hora_str):
    if not hora_str or str(hora_str).lower() == "none": return ""
    s_raw = str(hora_str).lower().strip()
    es_pm = "p.m." in s_raw or "pm" in s_raw or "p. m." in s_raw
    es_am = "a.m." in s_raw or "am" in s_raw or "a. m." in s_raw
    match = re.search(r'(\d{1,2}):(\d{2})', s_raw)
    if match:
        h, m = int(match.group(1)), match.group(2)
        if es_pm and h < 12: h += 12
        if es_am and h == 12: h = 0
        if h == 0 and es_pm: h = 12 
        if 8 <= h <= 11: suf = "de la mañana"
        elif h == 12: suf = "del día"
        elif 13 <= h <= 19: suf = "de la tarde"
        elif h >= 20: suf = "p.m."
        else: suf = "a.m."
        h_show = h - 12 if h > 12 else (12 if h == 0 else h)
        return f"{h_show}:{m} {suf}"
    return hora_str

def obtener_concat_texto(record):
    parts = []
    val_punto = record.get('Punto de reunion')
    val_ruta = record.get('Ruta a seguir')
    val_mun = record.get('Municipio')
    val_sec = record.get('Seccion')

    if val_punto and str(val_punto).lower() != 'none': parts.append(str(val_punto))
    if val_ruta and str(val_ruta).lower() != 'none': parts.append(str(val_ruta))
    if val_mun and str(val_mun).lower() != 'none': parts.append(f"Municipio {val_mun}")
    if val_sec and str(val_sec).lower() != 'none': parts.append(f"Sección {str(val_sec).upper()}")
    return ", ".join(parts)

# ============================================================
#  INICIO DE LA APP
# ============================================================

st.set_page_config(page_title="Provident Pro v129", layout="wide")

# 0. NAVEGACIÓN PERSISTENTE
if 'active_module' not in st.session_state:
    qp = st.query_params.get("view", "Calendario") 
    st.session_state.active_module = qp

def navegar_a(modulo):
    st.session_state.active_module = modulo
    st.query_params["view"] = modulo

# 1. BLOQUEO TECLADO
st.markdown("""
<script>
    const observer = new MutationObserver((mutations) => {
        const inputs = window.parent.document.querySelectorAll('input[type="text"]');
        inputs.forEach(input => {
            if (input.getAttribute('aria-autocomplete') === 'list') {
                input.setAttribute('inputmode', 'none');
                input.setAttribute('readonly', 'true');
            }
        });
    });
    observer.observe(window.parent.document.body, { childList: true, subtree: true });
</script>
""", unsafe_allow_html=True)

# 2. ESTILOS CSS - ÚNICAMENTE CALENDARIO
st.markdown("""
<style>
    /* CALENDARIO PERSONALIZADO (AZUL CORPORATIVO) */
    .cal-title {
        text-align: center; font-size: 1.5em; font-weight: bold; margin: 0 !important; padding-bottom: 10px; color: #333 !important; background-color: #fff;
    }
    .c-grid { display: grid; grid-template-columns: repeat(7, 1fr); gap: 2px; margin-top: 0px !important; }
    .c-head { background: #002060 !important; color: white !important; padding: 4px; text-align: center; font-weight: bold; border-radius: 2px; font-size: 14px; }
    .c-cell { background: white !important; border: 1px solid #ccc !important; border-radius: 2px; height: 160px; display: flex; flex-direction: column; justify-content: space-between; overflow: hidden; }
    .c-day { flex: 0 0 auto; background: #00b0f0 !important; color: white !important; font-weight: 900; font-size: 1.1em; text-align: center; padding: 2px 0; }
    .c-body { flex-grow: 1; width: 100%; background-position: center; background-repeat: no-repeat; background-size: cover; background-color: #f8f8f8 !important; }
    .c-foot { flex: 0 0 auto; height: 20px; background: #002060 !important; color: #ffffff !important; font-weight: 900; text-align: center; font-size: 0.9em; padding: 1px; white-space: nowrap; overflow: hidden; }
    .c-foot-empty { flex: 0 0 auto; height: 20px; background: #e0e0e0 !important; }
    
    @media (max-width: 600px) {
        .c-cell { height: 110px; }
        .c-day { font-size: 0.9em; }
        .c-foot, .c-foot-empty { font-size: 0.7em; height: 16px; }
    }
</style>
""", unsafe_allow_html=True)

if 'config' not in st.session_state:
    if os.path.exists("config_app.json"):
        with open("config_app.json", "r") as f: st.session_state.config = json.load(f)
    else: st.session_state.config = {"plantillas": {}}

st.title("🚀 Generador Pro v129")
TOKEN = "patyclv7hDjtGHB0F.19829008c5dee053cba18720d38c62ed86fa76ff0c87ad1f2d71bfe853ce9783"
headers = {"Authorization": f"Bearer {TOKEN}"}

# --- SIDEBAR ---
with st.sidebar:
    st.header("🔗 Conexión")
    r_bases = requests.get("https://api.airtable.com/v0/meta/bases", headers=headers)
    if r_bases.status_code == 200:
        base_opts = {b['name']: b['id'] for b in r_bases.json()['bases']}
        
        with st.expander("📂 Seleccionar Base", expanded=True):
            base_sel = st.radio("Bases disponibles:", list(base_opts.keys()), label_visibility="collapsed")
        
        if base_sel:
            st.session_state['base_activa_id'] = base_opts[base_sel]
            r_tab = requests.get(f"https://api.airtable.com/v0/meta/bases/{base_opts[base_sel]}/tables", headers=headers)
            if r_tab.status_code == 200:
                tablas_data = r_tab.json()['tables']
                st.session_state['todas_tablas'] = {t['name']: t['id'] for t in tablas_data}
                
                with st.expander("📅 Seleccionar Tabla (Mes)", expanded=True):
                    tabla_sel = st.radio("Tablas disponibles:", list(st.session_state['todas_tablas'].keys()), label_visibility="collapsed")
                
                if 'tabla_actual_nombre' not in st.session_state or st.session_state['tabla_actual_nombre'] != tabla_sel:
                    with st.spinner("Cargando datos automáticamente..."):
                        r_reg = requests.get(f"https://api.airtable.com/v0/{base_opts[base_sel]}/{st.session_state['todas_tablas'][tabla_sel]}", headers=headers)
                        recs = r_reg.json().get("records", [])
                        st.session_state.raw_data_original = recs
                        st.session_state.raw_records = [
                            {'id': r['id'], 'fields': {k: (procesar_texto_maestro(v, k) if k != 'Fecha' else v) for k, v in r['fields'].items()}}
                            for r in recs
                        ]
                        st.session_state['tabla_actual_nombre'] = tabla_sel
                    st.success(f"✅ {len(recs)} registros cargados")
                    st.rerun()

    st.divider()
    
    # --- MENÚ DE NAVEGACIÓN ---
    if 'raw_records' in st.session_state:
        st.subheader("⚡ Generar")
        
        t_pos = "primary" if st.session_state.active_module == "Postales" else "secondary"
        if st.button("📮 Postales", type=t_pos, use_container_width=True):
            navegar_a("Postales"); st.rerun()
            
        t_rep = "primary" if st.session_state.active_module == "Reportes" else "secondary"
        if st.button("📄 Reportes", type=t_rep, use_container_width=True):
            navegar_a("Reportes"); st.rerun()
            
        st.subheader("📅 Eventos")
        t_cal = "primary" if st.session_state.active_module == "Calendario" else "secondary"
        if st.button("📆 Calendario", type=t_cal, use_container_width=True):
            navegar_a("Calendario"); st.rerun()

        if st.button("💾 Guardar Config", use_container_width=True):
            with open("config_app.json", "w") as f: json.dump(st.session_state.config, f)
            st.toast("Guardado")

# --- MAIN ---
if 'raw_records' not in st.session_state:
    st.info("👈 Conecta una base.")
else:
    modulo = st.session_state.active_module
    df_full = pd.DataFrame([r['fields'] for r in st.session_state.raw_records])
    AZUL = RGBColor(0, 176, 240)

    # --------------------------------------------------------
    # MÓDULO POSTALES
    # --------------------------------------------------------
    if modulo == "Postales":
        st.subheader("📮 Generador de Postales")
        
        df_view = df_full.copy()
        for c in df_view.columns:
            if isinstance(df_view[c].iloc[0], list): df_view.drop(c, axis=1, inplace=True)
        
        # 1. CHECKBOX MAESTRO
        sel_all = st.checkbox("Seleccionar Todo", value=False, key="sel_all_postales")
        
        # 2. INSERTAR COLUMNA "✅" (VISUALMENTE LIMPIO)
        df_view.insert(0, "✅", sel_all)
        
        # 3. EDITOR DE DATOS
        df_edit = st.data_editor(df_view, hide_index=True)
        sel_idx = df_edit.index[df_edit["✅"]==True].tolist()
        
        if sel_idx:
            folder = os.path.join("Plantillas", "POSTALES")
            if not os.path.exists(folder): os.makedirs(folder)
            archs = [f for f in os.listdir(folder) if f.endswith('.pptx')]
            tipos = df_view.loc[sel_idx, "Tipo"].unique()
            cols = st.columns(len(tipos)) if len(tipos)>0 else [st]
            for i, t in enumerate(tipos):
                mem = st.session_state.config["plantillas"].get(t)
                idx = archs.index(mem) if mem in archs else 0
                st.session_state.config["plantillas"][t] = cols[i].selectbox(f"Plantilla '{t}':", archs, index=idx, key=f"pp_{t}")

            if st.button("🔥 GENERAR POSTALES", type="primary"):
                p_bar = st.progress(0); 
                archivos_generados = [] # LOCAL
                TAM_MAPA = {"<<Tipo>>":12, "<<Sucursal>>":44, "<<Seccion>>":12, "<<Conhora>>":32, "<<Concat>>":32, "<<Consuc>>":32, "<<Confechor>>":28, "<<Confecha>>":32}
                
                for i, idx in enumerate(sel_idx):
                    rec = st.session_state.raw_records[idx]['fields']
                    orig = st.session_state.raw_data_original[idx]['fields']
                    try: dt = datetime.strptime(rec.get('Fecha','2025-01-01'), '%Y-%m-%d'); nm = MESES_ES[dt.month - 1]
                    except: dt = datetime.now(); nm = "error"
                    
                    ft = rec.get('Tipo', 'Sin Tipo')
                    fs = rec.get('Sucursal', '000')
                    tfe = obtener_fecha_texto(dt); tho = obtener_hora_texto(rec.get('Hora',''))
                    
                    tfe_confechor = f"{nm.capitalize()} {dt.day} de {dt.year}"
                    fcf = f"{tfe_confechor.strip()}\n{tho.strip()}"
                    
                    fcc = f"Sucursal {fs}" if ft == "Actividad en Sucursal" else obtener_concat_texto(rec)
                    ftag = f"Sucursal {fs}" if ft == "Actividad en Sucursal" else fcc
                    narc = re.sub(r'[\\/*?:"<>|]', "", f"{dt.day} de {nm} de {dt.year} - {ft}, {fs} - {ftag}")[:120] + ".png"
                    
                    reps = {
                        "<<Tipo>>":textwrap.fill(ft,width=35), 
                        "<<Sucursal>>":fs, 
                        "<<Seccion>>":rec.get('Seccion'), 
                        "<<Confechor>>":fcf, 
                        "<<Concat>>":fcc, 
                        "<<Consuc>>":fcc,
                        "<<Confecha>>":tfe,
                        "<<Conhora>>":tho
                    }
                    
                    try: prs = Presentation(os.path.join(folder, st.session_state.config["plantillas"][ft]))
                    except: continue
                    if ft == "Actividad en Sucursal" and not orig.get("Lista de asistencia") and len(prs.slides)>=4: prs.slides._sldIdLst.remove(prs.slides._sldIdLst[3])

                    for slide in prs.slides:
                        for shp in slide.shapes:
                            if shp.has_text_frame:
                                for tag in ["Foto de equipo", "Foto 01", "Foto 02", "Foto 03", "Foto 04", "Foto 05", "Foto 06", "Foto 07", "Reporte firmado", "Lista de asistencia"]:
                                    if f"<<{tag}>>" in shp.text_frame.text:
                                        adj = orig.get(tag)
                                        if adj:
                                            try:
                                                io = procesar_imagen_inteligente(requests.get(adj[0]['url']).content, shp.width, shp.height, con_blur=True)
                                                slide.shapes.add_picture(io, shp.left, shp.top, shp.width, shp.height)
                                                shp._element.getparent().remove(shp._element)
                                            except: pass
                        for shp in slide.shapes:
                            if shp.has_text_frame:
                                for tag, val in reps.items():
                                    if tag in shp.text_frame.text:
                                        if tag in ["<<Tipo>>", "<<Sucursal>>"]: shp.top -= Pt(2)
                                        
                                        tf = shp.text_frame
                                        tf.vertical_anchor = MSO_ANCHOR.MIDDLE 
                                        tf.word_wrap = True
                                        
                                        if tag == "<<Confechor>>":
                                            tf.clear()
                                            p1 = tf.paragraphs[0]
                                            p1.text = tfe_confechor.strip()
                                            p1.alignment = PP_ALIGN.CENTER
                                            p1.font.bold = True
                                            p1.font.color.rgb = AZUL
                                            p1.font.size = Pt(28)
                                            
                                            p2 = tf.add_paragraph()
                                            p2.text = tho.strip()
                                            p2.alignment = PP_ALIGN.CENTER
                                            p2.font.bold = True
                                            p2.font.color.rgb = AZUL
                                            p2.font.size = Pt(28)
                                        else:
                                            bp = tf._element.bodyPr
                                            for c in ['spAutoFit', 'normAutofit', 'noAutofit']:
                                                if bp.find(qn(f'a:{c}')) is not None: bp.remove(bp.find(qn(f'a:{c}')))
                                            if tag=="<<Tipo>>": bp.append(tf._element.makeelement(qn('a:spAutoFit')))
                                            else: bp.append(tf._element.makeelement(qn('a:normAutofit')))
                                            
                                            tf.clear(); p = tf.paragraphs[0]
                                            if tag in ["<<Confecha>>", "<<Conhora>>", "<<Consuc>>"]: p.alignment = PP_ALIGN.CENTER
                                            
                                            p.space_before=Pt(0); p.space_after=Pt(0); p.line_spacing=1.0
                                            run = p.add_run(); run.text=str(val); run.font.bold=True; run.font.color.rgb=AZUL
                                            run.font.size=Pt(TAM_MAPA.get(tag,12))
                    
                    pp_io = BytesIO(); prs.save(pp_io)
                    dout = generar_pdf(pp_io.getvalue())
                    if dout:
                        imgs = convert_from_bytes(dout, dpi=170, fmt='jpeg')
                        with BytesIO() as b: 
                            imgs[0].save(b, format="JPEG", quality=85, optimize=True, progressive=True)
                            fbytes = b.getvalue()
                        path = f"{dt.year}/{str(dt.month).zfill(2)} - {nm}/Postales/{fs}/{narc}"
                        archivos_generados.append({"RutaZip":path, "Datos":fbytes})
                    p_bar.progress((i+1)/len(sel_idx))
                
                # BOTÓN DE DESCARGA INMEDIATO
                if archivos_generados:
                    zb = BytesIO()
                    with zipfile.ZipFile(zb, "a", zipfile.ZIP_DEFLATED) as z:
                        for f in archivos_generados: z.writestr(f["RutaZip"], f["Datos"])
                    st.success("✅ Generación Completada")
                    st.download_button(f"⬇️ DESCARGAR {len(archivos_generados)} POSTALES", zb.getvalue(), f"Postales_{datetime.now().strftime('%H%M%S')}.zip", "application/zip", type="primary")

    # --------------------------------------------------------
    # MÓDULO REPORTES
    # --------------------------------------------------------
    elif modulo == "Reportes":
        st.subheader("📄 Generador de Reportes")
        df_view = df_full.copy()
        for c in df_view.columns:
            if isinstance(df_view[c].iloc[0], list): df_view.drop(c, axis=1, inplace=True)
        
        # 1. CHECKBOX MAESTRO
        sel_all_r = st.checkbox("Seleccionar Todo", value=False, key="sel_all_reportes")
        
        # 2. INSERTAR COLUMNA "✅"
        df_view.insert(0, "✅", sel_all_r)
        
        # 3. EDITOR
        df_edit = st.data_editor(df_view, hide_index=True, key="ed_r")
        sel_idx = df_edit.index[df_edit["✅"]==True].tolist()
        
        if sel_idx:
            folder = os.path.join("Plantillas", "REPORTES")
            if not os.path.exists(folder): os.makedirs(folder)
            archs = [f for f in os.listdir(folder) if f.endswith('.pptx')]
            tipos = df_view.loc[sel_idx, "Tipo"].unique()
            cols = st.columns(len(tipos)) if len(tipos)>0 else [st]
            for i, t in enumerate(tipos):
                mem = st.session_state.config["plantillas"].get(t)
                idx = archs.index(mem) if mem in archs else 0
                st.session_state.config["plantillas"][t] = cols[i].selectbox(f"Plantilla '{t}':", archs, index=idx, key=f"pr_{t}")

            if st.button("🔥 CREAR REPORTES", type="primary"):
                p_bar = st.progress(0)
                archivos_generados = [] # LOCAL
                TAM_MAPA = {"<<Tipo>>":12, "<<Sucursal>>":12, "<<Seccion>>":12, "<<Confecha>>":24, "<<Conhora>>":15, "<<Consuc>>":24, "<<Confechor>>":20, "<<Concat>>":12}
                
                for i, idx in enumerate(sel_idx):
                    rec = st.session_state.raw_records[idx]['fields']
                    orig = st.session_state.raw_data_original[idx]['fields']
                    try: dt = datetime.strptime(rec.get('Fecha','2025-01-01'), '%Y-%m-%d'); nm = MESES_ES[dt.month - 1]
                    except: dt = datetime.now(); nm = "error"
                    
                    ft = rec.get('Tipo', 'Sin Tipo'); fs = rec.get('Sucursal', '000')
                    tfe = obtener_fecha_texto(dt); tho = obtener_hora_texto(rec.get('Hora',''))
                    
                    tfe_confechor = f"{nm.capitalize()} {dt.day} de {dt.year}"
                    fcf = f"{tfe_confechor.strip()}\n{tho.strip()}"
                    
                    fcs = f"Sucursal {fs}" if ft == "Actividad en Sucursal" else ""
                    fcc = f"Sucursal {fs}" if ft == "Actividad en Sucursal" else obtener_concat_texto(rec)
                    ftag = fcs if ft == "Actividad en Sucursal" else obtener_concat_texto(rec)
                    narc = re.sub(r'[\\/*?:"<>|]', "", f"{dt.day} de {nm} de {dt.year} - {ft}, {fs} - {ftag}")[:120] + ".pdf"
                    
                    reps = {
                        "<<Tipo>>":textwrap.fill(ft,width=35), 
                        "<<Sucursal>>":fs, 
                        "<<Seccion>>":rec.get('Seccion'), 
                        "<<Confecha>>":tfe, 
                        "<<Conhora>>":tho, 
                        "<<Consuc>>":fcs,
                        "<<Confechor>>":fcf,
                        "<<Concat>>":fcc
                    }
                    
                    try: prs = Presentation(os.path.join(folder, st.session_state.config["plantillas"][ft]))
                    except: continue
                    if ft == "Actividad en Sucursal" and not orig.get("Lista de asistencia") and len(prs.slides)>=4: prs.slides._sldIdLst.remove(prs.slides._sldIdLst[3])

                    for slide in prs.slides:
                        for shp in slide.shapes:
                            if shp.has_text_frame:
                                for tag in ["Foto de equipo", "Foto 01", "Foto 02", "Foto 03", "Foto 04", "Foto 05", "Foto 06", "Foto 07", "Reporte firmado", "Lista de asistencia"]:
                                    if f"<<{tag}>>" in shp.text_frame.text:
                                        adj = orig.get(tag)
                                        if adj:
                                            try:
                                                io = procesar_imagen_inteligente(requests.get(adj[0]['url']).content, shp.width, shp.height, con_blur=True)
                                                slide.shapes.add_picture(io, shp.left, shp.top, shp.width, shp.height)
                                                shp._element.getparent().remove(shp._element)
                                            except: pass
                        for shp in slide.shapes:
                            if shp.has_text_frame:
                                for tag, val in reps.items():
                                    if tag in shp.text_frame.text:
                                        if tag in ["<<Tipo>>", "<<Sucursal>>"]: shp.top -= Pt(2)
                                        
                                        tf = shp.text_frame
                                        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                                        tf.word_wrap = True
                                        
                                        if tag == "<<Confechor>>":
                                            tf.clear()
                                            p1 = tf.paragraphs[0]
                                            p1.text = tfe_confechor.strip()
                                            p1.alignment = PP_ALIGN.CENTER
                                            p1.font.bold = True
                                            p1.font.color.rgb = AZUL
                                            p1.font.size = Pt(20)
                                            
                                            p2 = tf.add_paragraph()
                                            p2.text = tho.strip()
                                            p2.alignment = PP_ALIGN.CENTER
                                            p2.font.bold = True
                                            p2.font.color.rgb = AZUL
                                            p2.font.size = Pt(20)
                                        else:
                                            bp = tf._element.bodyPr
                                            for c in ['spAutoFit', 'normAutofit', 'noAutofit']:
                                                if bp.find(qn(f'a:{c}')) is not None: bp.remove(bp.find(qn(f'a:{c}')))
                                            if tag=="<<Tipo>>": bp.append(tf._element.makeelement(qn('a:spAutoFit')))
                                            elif tag=="<<Conhora>>": bp.append(tf._element.makeelement(qn('a:spAutoFit')))
                                            else: bp.append(tf._element.makeelement(qn('a:normAutofit')))
                                            
                                            tf.clear(); p = tf.paragraphs[0]
                                            if tag in ["<<Confecha>>", "<<Conhora>>", "<<Consuc>>"]: p.alignment = PP_ALIGN.CENTER
                                            
                                            p.space_before=Pt(0); p.space_after=Pt(0); p.line_spacing=1.0
                                            run = p.add_run(); run.text=str(val); run.font.bold=True; run.font.color.rgb=AZUL
                                            run.font.size=Pt(TAM_MAPA.get(tag,12))
                    
                    pp_io = BytesIO(); prs.save(pp_io)
                    dout = generar_pdf(pp_io.getvalue())
                    if dout:
                        path = f"{dt.year}/{str(dt.month).zfill(2)} - {nm}/Reportes/{fs}/{narc}"
                        archivos_generados.append({"RutaZip":path, "Datos":dout})
                    p_bar.progress((i+1)/len(sel_idx))
                
                # BOTÓN DE DESCARGA INMEDIATO
                if archivos_generados:
                    zb = BytesIO()
                    with zipfile.ZipFile(zb, "a", zipfile.ZIP_DEFLATED) as z:
                        for f in archivos_generados: z.writestr(f["RutaZip"], f["Datos"])
                    st.success("✅ Generación Completada")
                    st.download_button(f"⬇️ DESCARGAR {len(archivos_generados)} REPORTES", zb.getvalue(), f"Reportes_{datetime.now().strftime('%H%M%S')}.zip", "application/zip", type="primary")

    # --------------------------------------------------------
    # MÓDULO CALENDARIO
    # --------------------------------------------------------
    elif modulo == "Calendario":
        st.subheader("📅 Calendario de Actividades")
        
        if 'todas_tablas' in st.session_state:
            nombres_tablas = list(st.session_state['todas_tablas'].keys())
            idx_actual = 0
            if 'tabla_actual_nombre' in st.session_state and st.session_state['tabla_actual_nombre'] in nombres_tablas:
                idx_actual = nombres_tablas.index(st.session_state['tabla_actual_nombre'])
            
            with st.expander("📅 Cambiar Mes (Tabla)", expanded=False):
                # CARGA AUTOMATICA
                nueva_tabla = st.radio("Meses:", nombres_tablas, index=idx_actual, horizontal=True)
            
            if nueva_tabla != st.session_state.get('tabla_actual_nombre'):
                with st.spinner(f"Cargando {nueva_tabla}..."):
                    base_id = st.session_state['base_activa_id']
                    table_id = st.session_state['todas_tablas'][nueva_tabla]
                    r_reg = requests.get(f"https://api.airtable.com/v0/{base_id}/{table_id}", headers=headers)
                    recs = r_reg.json().get("records", [])
                    st.session_state.raw_data_original = recs
                    st.session_state.raw_records = [
                        {'id': r['id'], 'fields': {k: (procesar_texto_maestro(v, k) if k != 'Fecha' else v) for k, v in r['fields'].items()}}
                        for r in recs
                    ]
                    st.session_state['tabla_actual_nombre'] = nueva_tabla
                st.rerun()

        st.divider()

        fechas_oc = {}
        fechas_lista = []
        for r in st.session_state.raw_data_original:
            f = r['fields'].get('Fecha')
            if f:
                f_short = f.split('T')[0]
                if f_short not in fechas_oc: fechas_oc[f_short] = []
                th = None
                if 'Postal' in r['fields']:
                    att = r['fields']['Postal']
                    if isinstance(att, list) and len(att)>0: 
                        th = att[0].get('url') 
                fechas_oc[f_short].append({"id":r['id'], "thumb":th})
                fechas_lista.append(f_short)

        if not fechas_oc:
            st.warning("No hay fechas en esta tabla.")
        else:
            dt_objs = [datetime.strptime(x, '%Y-%m-%d') for x in fechas_lista]
            mc = Counter([(d.year, d.month) for d in dt_objs])
            ay, am = mc.most_common(1)[0][0]
            
            cal = calendar.Calendar(firstweekday=0) 
            weeks = cal.monthdayscalendar(ay, am)
            
            st.markdown("""
            <style>
            .cal-title {
                text-align: center; font-size: 1.5em; font-weight: bold; margin: 0 !important; padding-bottom: 10px; color: #333 !important; background-color: #fff;
            }
            .c-grid { display: grid; grid-template-columns: repeat(7, 1fr); gap: 2px; margin-top: 0px !important; }
            .c-head { background: #002060 !important; color: white !important; padding: 4px; text-align: center; font-weight: bold; border-radius: 2px; font-size: 14px; }
            .c-cell { background: white !important; border: 1px solid #ccc !important; border-radius: 2px; height: 160px; display: flex; flex-direction: column; justify-content: space-between; overflow: hidden; }
            .c-day { flex: 0 0 auto; background: #00b0f0 !important; color: white !important; font-weight: 900; font-size: 1.1em; text-align: center; padding: 2px 0; }
            .c-body { flex-grow: 1; width: 100%; background-position: center; background-repeat: no-repeat; background-size: cover; background-color: #f8f8f8 !important; }
            .c-foot { flex: 0 0 auto; height: 20px; background: #002060 !important; color: #ffffff !important; font-weight: 900; text-align: center; font-size: 0.9em; padding: 1px; white-space: nowrap; overflow: hidden; }
            .c-foot-empty { flex: 0 0 auto; height: 20px; background: #e0e0e0 !important; }
            
            @media (max-width: 600px) {
                .c-cell { height: 110px; }
                .c-day { font-size: 0.9em; }
                .c-foot, .c-foot-empty { font-size: 0.7em; height: 16px; }
            }
            </style>
            """, unsafe_allow_html=True)

            h = f"<div class='cal-title'>📅 {MESES_ES[am-1].capitalize()} {ay}</div>"
            h += "<div class='c-grid'>"
            for d in ["LUN","MAR","MIÉ","JUE","VIE","SÁB","DOM"]: h += f"<div class='c-head'>{d}</div>"
            
            for wk in weeks:
                for d in wk:
                    if d == 0: h += "<div class='c-cell' style='border:none; background:transparent;'></div>"
                    else:
                        k = f"{ay}-{str(am).zfill(2)}-{str(d).zfill(2)}"
                        acts = fechas_oc.get(k, [])
                        
                        h += f"<div class='c-cell'>"
                        
                        # HEADER
                        h += f"<div class='c-day'>{d}</div>"
                        
                        # BODY
                        style_bg = ""
                        if acts and acts[0]['thumb']:
                            style_bg = f"style=\"background-image: url('{acts[0]['thumb']}');\""
                        h += f"<div class='c-body' {style_bg}></div>"
                        
                        # FOOTER LOGIC
                        if len(acts) > 1:
                            h += f"<div class='c-foot'>+ {len(acts)-1} más</div>"
                        elif len(acts) == 1:
                            h += f"<div class='c-foot'>&nbsp;</div>"
                        else:
                            h += f"<div class='c-foot-empty'></div>"
                        
                        h += "</div>"
            h += "</div>"
            st.markdown(h, unsafe_allow_html=True)
